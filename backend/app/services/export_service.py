import os
from datetime import datetime

from docx import Document
from pptx import Presentation


EXPORT_DIR = os.path.join(os.path.dirname(__file__), "exports")

# ✅ ensure export folder exists
os.makedirs(EXPORT_DIR, exist_ok=True)


# ---------------------------------------------------
# ✅ EXPORT DOCX
# ---------------------------------------------------
def export_to_docx(project_title: str, sections: dict) -> str:
    """
    sections = {
        "Introduction": "content...",
        "Problem Statement": "content...",
    }
    """
    doc = Document()
    doc.add_heading(project_title, level=1)

    for title, content in sections.items():
        doc.add_heading(title, level=2)
        doc.add_paragraph(content)

    filename = f"{project_title}_{timestamp()}.docx"
    filepath = os.path.join(EXPORT_DIR, filename)
    doc.save(filepath)

    return filepath


# ---------------------------------------------------
# ✅ EXPORT PPTX
# ---------------------------------------------------
def export_to_pptx(project_title: str, sections: dict) -> str:
    pres = Presentation()

    # ✅ title slide
    title_slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project_title

    for title, content in sections.items():
        slide_layout = pres.slide_layouts[1]  # title + content layout
        slide = pres.slides.add_slide(slide_layout)
        slide.shapes.title.text = title

        textbox = slide.shapes.placeholders[1].text_frame
        textbox.text = content

    filename = f"{project_title}_{timestamp()}.pptx"
    filepath = os.path.join(EXPORT_DIR, filename)
    pres.save(filepath)

    return filepath


# ---------------------------------------------------
# ✅ Helper
# ---------------------------------------------------
def timestamp():
    return datetime.now().strftime("%Y%m%d_%H%M%S")
